#import pymupdf  # PyMuPDF
import sys
import os
import subprocess
import unicodedata
import pymupdf.layout
import pymupdf4llm
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
import html2text_rs
import pandas as pd
from utils import log_it

# ------------------------------
# requirements:
# ------------------------------
# pip install pymupdf
# pip install pymupdf-layout
# pip install python-docx
# pip install pandas openpyxl
# pip install xlrd
# pip install python-pptx
# pip install html2text_rs
# INSTALL libreoffice 


# -----------------------
class DocConverter:
# -----------------------

    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def __init__(self):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        self.unusual_to_usual_char = {
            '\u2010' : '-', 
            '\u2011' : '-',
            '\u2012' : '-',
            '\u2013' : '-',
            '\u2014' : '-',
            '\u2015' : '-',
            '\u0142' : 'l',
            'ø'      : 'o',
            'Ø'      : 'O',
            'æ'      : 'ae',
            'Ł'      : 'L',
        }
        # Mapping dictionary for Greek letters to their English names
        self.greek_to_english = {
            'Α': 'Alpha', 'α': 'alpha',
            'Β': 'Beta', 'β': 'beta',
            'Γ': 'Gamma', 'γ': 'gamma',
            'Δ': 'Delta', 'δ': 'delta',
            'Ε': 'Epsilon', 'ε': 'epsilon',
            'Ζ': 'Zeta', 'ζ': 'zeta',
            'Η': 'Eta', 'η': 'eta',
            'Θ': 'Theta', 'θ': 'theta',
            'Ι': 'Iota', 'ι': 'iota',
            'Κ': 'Kappa', 'κ': 'kappa',
            'Λ': 'Lambda', 'λ': 'lambda',
            'Μ': 'Mu', 'μ': 'mu',
            'Ν': 'Nu', 'ν': 'nu',
            'Ξ': 'Xi', 'ξ': 'xi',
            'Ο': 'Omicron', 'ο': 'omicron',
            'Π': 'Pi', 'π': 'pi',
            'Ρ': 'Rho', 'ρ': 'rho',
            'Σ': 'Sigma', 'σ': 'sigma', 'ς': 'sigma',
            'Τ': 'Tau', 'τ': 'tau',
            'Υ': 'Upsilon', 'υ': 'upsilon',
            'Φ': 'Phi', 'φ': 'phi',
            'Χ': 'Chi', 'χ': 'chi',
            'Ψ': 'Psi', 'ψ': 'psi',
            'Ω': 'Omega', 'ω': 'omega'
        }


    # - - - - - - - - - - - - - - - - - - - - - - - - - - 
    def replace_misleading_chars(self, text):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - 
        new_chars = list()
        for char in text: new_chars.append(self.unusual_to_usual_char.get(char, char))
        return ''.join(new_chars)
    

    # - - - - - - - - - - - - - - - - - - - - - - - - - - 
    def cello_normalized(self, input_str):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - 
        if not input_str: return None
        tmp1 = self.replace_misleading_chars(input_str)
        #if tmp1 != input_str: print(f"INFO replaced some unusual character at step 1 in '{input_str}'")
        tmp2 = self.remove_accents(tmp1)
        #if tmp2 != tmp1: print(f"INFO replaced some unusual character at step 2 in '{input_str}'")
        tmp3 = self.translate_greek_to_english_names(tmp2)
        #if tmp3 != tmp2: print(f"INFO replaced some unusual character at step 3 in '{input_str}'")
        return tmp3


    # - - - - - - - - - - - - - - - - - - - - - - - - - - 
    def remove_accents(self, input_str):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - 
        nfkd_form = unicodedata.normalize('NFKD', input_str)
        return u"".join([c for c in nfkd_form if not unicodedata.combining(c)])


    # - - - - - - - - - - - - - - - - - - - - - - - - - - 
    def translate_greek_to_english_names(self, text):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - 
        new_chars = list()
        for char in text: new_chars.append(self.greek_to_english.get(char, char))
        return ''.join(new_chars)

    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def save_text_file(self, text_content, filename):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        if not text_content: # None or empty string
            log_it("ERROR Converted file is empty:", filename)
            return
        log_it("INFO Saving", filename)
        with open(filename, "w", encoding="utf-8", errors="replace") as f:
            f.write(text_content)


    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def convert_to_text(self, input_file, publication_id=None):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        log_it("INFO Converting file", input_file)
        right_most_dot_pos = input_file.rfind(".")
        if right_most_dot_pos == -1: return None
        ext = input_file[right_most_dot_pos:]
        raw_text = None
        if ext == ".pdf": raw_text = self.convert_pdf_to_text(input_file, publication_id=publication_id)
        elif ext == ".doc": raw_text = self.convert_doc_to_text(input_file)
        elif ext == ".docx": raw_text = self.convert_docx_to_text(input_file)
        elif ext == ".xlsx": raw_text = self.convert_excel_to_text(input_file)
        elif ext == ".xls": raw_text = self.convert_excel_to_text(input_file)
        elif ext == ".pptx": raw_text = self.convert_pptx_to_text(input_file)
        elif ext == ".html": raw_text = self.convert_html_to_text(input_file)
        return self.cello_normalized(raw_text)


    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def convert_html_to_text(self, input_file):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        file_name = f"\n\n## HTML file {os.path.basename(input_file)}\n\n" # md-like title for file
        html = open(input_file).read()
        #text_string = html2text_rs.text_markdown(html)
        text_string = html2text_rs.text_plain(html) 
        return "\n".join([file_name, text_string])


    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def convert_pptx_to_text(self, input_file):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        prs = Presentation(input_file)
        prs_name = f"\n\n## PowerPoint presentation file {os.path.basename(input_file)}\n\n" # md-like title for file
        text_runs = [prs_name]
        for i, slide in enumerate(prs.slides):
            text_runs.append(f"\n\n### Slide {i+1}\n\n") # md-like title for each slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)


    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def convert_excel_to_text(self, input_file):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        log_it("input file", input_file)
        all_sheets = pd.read_excel(input_file, sheet_name=None)
        
        text_list = list()
        xls_name = f"\n\n## Excel file {os.path.basename(input_file)}\n\n" # md-like title for file
        text_list.append(xls_name)
        for sheet_name, df in all_sheets.items():
            sheet_title = f"\n\n### Excel sheet {sheet_name}\n\n" # md-like title for each sheet
            text_list.append(sheet_title)
            sheet_data = df.to_csv(sep='\t', index=False, header=True)
            text_list.append(sheet_data)
        return "\n".join(text_list) 


    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def convert_pdf_to_text(self, input_file, publication_id=None, max_pages=25):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        debug = False
        ori_doc = pymupdf.open(input_file)
        num_pages = min(len(ori_doc), max_pages)

        if debug: log_it("doc.name", ori_doc.name)
        if debug: log_it(f"Total pages in document: {num_pages}")
        num_pages = len(ori_doc)
        if num_pages > max_pages: 
            log_it(f"WARNING Ignoring pages after page {max_pages}: {input_file}")
            num_pages = max_pages
        new_doc = pymupdf.open() # new empty doc
        for i in range(num_pages):
            page = ori_doc[i]
            text = page.get_text()
            if text.strip():
                if debug: log_it(f"Inserting page {i+1}")
                new_doc.insert_pdf(ori_doc, from_page=i, to_page=i)
            else:
                if debug: log_it(f"Page {i+1} not inserted, is empty")
        try:
            output = pymupdf4llm.to_markdown(new_doc, footer=False, header=False)    # structure better for LLM (close to tsv)
            #output = pymupdf4llm.to_text(new_doc, footer=False, header=False)       # display nicer for human (but harder to parse)
        except Exception as e:
            log_it(f"ERROR during conversion of {input_file}", e)
            output = ""
        if publication_id is None: 
            return output
        else:
            return "\n\n".join([publication_id, output])


    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def convert_doc_to_text(self, input_file):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        command = ["libreoffice", "--headless", "--convert-to", "docx", input_file]
        try:
            subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            tmp_file = os.path.basename(input_file) + "x" # libreoffice writes in curr dir and change extension to .docx
            output = self.convert_docx_to_text(tmp_file)
            if os.path.exists(tmp_file): os.remove(tmp_file)
            return output
        except subprocess.CalledProcessError as e:
            log_it(f"ERROR during conversion of {input_file} to docx", e.stderr.decode())
            return None


    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
    def convert_docx_to_text(self, input_file):
    # - - - - - - - - - - - - - - - - - - - - - - - - - - - - - -
        def iter_block_items(parent):
            for child in parent.element.body.iterchildren():
                if child.tag.endswith('}p'):
                    yield Paragraph(child, parent)
                elif child.tag.endswith('}tbl'):
                    yield Table(child, parent)
        doc = Document(input_file)
        doc_name = f"\n\n## Document file {os.path.basename(input_file)}\n\n" # md-like title for file
        full_text_parts = [doc_name]
        for block in iter_block_items(doc):
            if isinstance(block, Paragraph):
                full_text_parts.append(block.text)
            elif isinstance(block, Table):
                table_text = []
                for row in block.rows:
                    row_text = "\t".join(cell.text for cell in row.cells)   # tsv is cool for LLM input
                    table_text.append(row_text)
                full_text_parts.append("\n".join(table_text))
        return "\n\n".join(full_text_parts)




# =============================================================
if __name__ == '__main__':
# =============================================================

    # some tests below

    input_dir="pdf"
    output_dir = "txt_data"
    converter = DocConverter()

    for file in ["8412300.pdf", "38565739_supp_material.docx", "29750960_supp_material.doc", "29627726_supp_material.pptx", 
                   "10741968_supp_material.html", "29750960_supp_material.doc", "38180245_supp_material_1.xlsx", "38358347_supp_material_2.xlsx", 
                   "31733513_supp_material.xlsx", "24662767_supp_material.xls",   "25984343_supp_material.xls", "31395879_supp_material_2.xls"]:
        input_file = f"{input_dir}/{file}"
        output = converter.convert_to_text(input_file)
        right_most_dot_pos = file.rfind(".")
        file_name = file[0:right_most_dot_pos] + ".txt"
        output_file = f"{output_dir}/{file_name}"
        converter.save_text_file(output, output_file)

